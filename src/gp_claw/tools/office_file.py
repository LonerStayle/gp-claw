import csv
import io
from pathlib import Path

from langchain_core.tools import tool

from gp_claw.security import validate_path


def create_office_tools(workspace_root: str) -> list:
    """사무용 파일 도구 생성. 실행 전 사용자 승인 필수."""

    @tool
    def excel_write(path: str, sheets: list[dict]) -> dict:
        """엑셀(.xlsx) 파일을 생성합니다. (승인 필요)

        Args:
            path: 워크스페이스 내 파일 경로 (.xlsx)
            sheets: 시트 목록. 각 시트는 {"name": "시트명", "headers": ["열1", ...], "rows": [["값1", ...], ...]} 형태
        """
        import openpyxl

        validated = validate_path(path, workspace_root)
        validated.parent.mkdir(parents=True, exist_ok=True)

        wb = openpyxl.Workbook()
        # 기본 시트 제거
        wb.remove(wb.active)

        for sheet_data in sheets:
            ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))
            headers = sheet_data.get("headers", [])
            if headers:
                ws.append(headers)
            for row in sheet_data.get("rows", []):
                ws.append(row)

        wb.save(str(validated))
        return {
            "path": str(validated),
            "size_bytes": validated.stat().st_size,
            "action": "created",
            "sheets": len(sheets),
        }

    @tool
    def csv_write(path: str, headers: list[str], rows: list[list]) -> dict:
        """CSV 파일을 생성합니다. (승인 필요)

        Args:
            path: 워크스페이스 내 파일 경로 (.csv)
            headers: 열 이름 목록
            rows: 데이터 행 목록
        """
        validated = validate_path(path, workspace_root)
        validated.parent.mkdir(parents=True, exist_ok=True)

        with open(validated, "w", newline="", encoding="utf-8-sig") as f:
            writer = csv.writer(f)
            writer.writerow(headers)
            writer.writerows(rows)

        return {
            "path": str(validated),
            "size_bytes": validated.stat().st_size,
            "action": "created",
            "rows": len(rows),
        }

    @tool
    def pdf_write(path: str, title: str, content: str) -> dict:
        """PDF 파일을 생성합니다. (승인 필요)

        Args:
            path: 워크스페이스 내 파일 경로 (.pdf)
            title: PDF 제목
            content: 본문 내용 (줄바꿈으로 단락 구분)
        """
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        validated = validate_path(path, workspace_root)
        validated.parent.mkdir(parents=True, exist_ok=True)

        # 한글 폰트 등록 시도
        font_name = "Helvetica"
        korean_font_paths = [
            "/System/Library/Fonts/AppleSDGothicNeo.ttc",
            "/usr/share/fonts/truetype/noto/NotoSansKR-Regular.ttf",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        ]
        for fp in korean_font_paths:
            if Path(fp).exists():
                try:
                    pdfmetrics.registerFont(TTFont("KoreanFont", fp))
                    font_name = "KoreanFont"
                    break
                except Exception:
                    continue

        doc = SimpleDocTemplate(
            str(validated),
            pagesize=A4,
            leftMargin=20 * mm,
            rightMargin=20 * mm,
            topMargin=20 * mm,
            bottomMargin=20 * mm,
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "CustomTitle", parent=styles["Title"], fontName=font_name, fontSize=18
        )
        body_style = ParagraphStyle(
            "CustomBody", parent=styles["Normal"], fontName=font_name, fontSize=11, leading=16
        )

        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 10 * mm))

        for paragraph in content.split("\n"):
            paragraph = paragraph.strip()
            if paragraph:
                story.append(Paragraph(paragraph, body_style))
                story.append(Spacer(1, 3 * mm))

        doc.build(story)
        return {
            "path": str(validated),
            "size_bytes": validated.stat().st_size,
            "action": "created",
        }

    @tool
    def pptx_write(path: str, title: str, slides: list[dict]) -> dict:
        """파워포인트(.pptx) 파일을 생성합니다. (승인 필요)

        Args:
            path: 워크스페이스 내 파일 경로 (.pptx)
            title: 발표 제목 (첫 번째 타이틀 슬라이드에 사용)
            slides: 슬라이드 목록. 각 슬라이드는 {"title": "제목", "content": "내용"} 형태
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt

        validated = validate_path(path, workspace_root)
        validated.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()

        # 타이틀 슬라이드
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = ""

        # 컨텐츠 슬라이드
        content_layout = prs.slide_layouts[1]
        for slide_data in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            body.text = slide_data.get("content", "")

        prs.save(str(validated))
        return {
            "path": str(validated),
            "size_bytes": validated.stat().st_size,
            "action": "created",
            "slides": len(slides) + 1,
        }

    return [excel_write, csv_write, pdf_write, pptx_write]
